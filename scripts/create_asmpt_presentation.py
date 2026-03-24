"""
Generate ASMPT progress update presentation (23-03-2026).
Run with: conda run -n GraduationProject python scripts/create_asmpt_presentation.py
Output: presentations/asmpt-23-03-2026.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Colours (TU/e-inspired) ─────────────────────────────────────────────────
BLUE_DARK   = RGBColor(0x00, 0x2B, 0x5C)   # TU/e dark blue
BLUE_MID    = RGBColor(0x00, 0x56, 0xA2)   # accent blue
ORANGE      = RGBColor(0xF0, 0x7D, 0x00)   # TU/e orange
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY  = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GREY   = RGBColor(0x33, 0x33, 0x33)
GREEN_DARK  = RGBColor(0x1A, 0x7A, 0x3C)

# ── Slide dimensions (16:9 widescreen) ──────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK_LAYOUT = prs.slide_layouts[6]   # completely blank


# ── Helper functions ────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=DARK_GREY, align=PP_ALIGN.LEFT,
                italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_title_bar(slide, title_text, subtitle_text=None):
    """Dark blue header bar with title."""
    bar = add_rect(slide, 0, 0, W, Inches(1.3), fill_color=BLUE_DARK)
    # Title
    add_textbox(slide, Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.7),
                title_text, font_size=28, bold=True, color=WHITE)
    if subtitle_text:
        add_textbox(slide, Inches(0.4), Inches(0.82), Inches(12.5), Inches(0.4),
                    subtitle_text, font_size=14, color=ORANGE)
    # Orange accent line below bar
    add_rect(slide, 0, Inches(1.3), W, Inches(0.06), fill_color=ORANGE)


def add_bullet_block(slide, left, top, width, height,
                     items, font_size=16, indent=None):
    """
    items: list of (level, text) tuples.  level 0 = bullet, level 1 = sub-bullet.
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    first = True
    for level, text in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.level = level
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = DARK_GREY
        if level == 0:
            p.space_before = Pt(4)
        # bullet symbol
        if level == 0:
            p.text = "• " + text
            run.text = "• " + text
        else:
            p.text = "    ◦ " + text
            run.text = "    ◦ " + text


def add_footer(slide, date="23 March 2026", page=None):
    add_textbox(slide, Inches(0.3), Inches(7.1), Inches(8), Inches(0.35),
                f"Graduation Project — ASMPT Update    |    {date}",
                font_size=9, color=RGBColor(0x88, 0x88, 0x88))
    if page:
        add_textbox(slide, Inches(12.5), Inches(7.1), Inches(0.8), Inches(0.35),
                    str(page), font_size=9, color=RGBColor(0x88, 0x88, 0x88),
                    align=PP_ALIGN.RIGHT)


def section_box(slide, left, top, width, height, title, color=BLUE_MID):
    """Coloured section header rectangle."""
    add_rect(slide, left, top, width, Inches(0.38), fill_color=color)
    add_textbox(slide, left + Inches(0.1), top + Inches(0.03),
                width - Inches(0.2), Inches(0.32),
                title, font_size=13, bold=True, color=WHITE)
    return top + Inches(0.38)


def status_badge(slide, left, top, text, done=True):
    color = GREEN_DARK if done else ORANGE
    add_rect(slide, left, top, Inches(1.3), Inches(0.32), fill_color=color)
    add_textbox(slide, left + Inches(0.05), top + Inches(0.03),
                Inches(1.2), Inches(0.26),
                text, font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = WHITE

# Large blue block top half
add_rect(slide, 0, 0, W, Inches(4.2), fill_color=BLUE_DARK)

# Orange accent stripe
add_rect(slide, 0, Inches(4.2), W, Inches(0.1), fill_color=ORANGE)

# Title text
add_textbox(slide, Inches(0.6), Inches(0.8), Inches(12.0), Inches(1.2),
            "LPV-LFR Model Augmentation",
            font_size=40, bold=True, color=WHITE)
add_textbox(slide, Inches(0.6), Inches(2.0), Inches(12.0), Inches(0.8),
            "Progress Update for ASMPT",
            font_size=28, bold=False, color=ORANGE)
add_textbox(slide, Inches(0.6), Inches(2.9), Inches(12.0), Inches(0.6),
            "Dual-Gantry First-Principles Model → Data-Augmented LPV System",
            font_size=18, color=RGBColor(0xBB, 0xCC, 0xDD))

# Lower content area
add_textbox(slide, Inches(0.6), Inches(4.6), Inches(6), Inches(0.4),
            "Graduation Project — TU Eindhoven", font_size=14, color=DARK_GREY)
add_textbox(slide, Inches(0.6), Inches(5.0), Inches(6), Inches(0.4),
            "23 March 2026", font_size=14, color=DARK_GREY)
add_textbox(slide, Inches(0.6), Inches(5.6), Inches(9), Inches(0.5),
            "Supervisors: A. Mesbah  |  Industry partner: ASMPT",
            font_size=12, color=RGBColor(0x77, 0x77, 0x77))


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = WHITE
add_title_bar(slide, "Agenda")
add_footer(slide, page=2)

items = [
    (0, "1.   Project context & goal"),
    (0, "2.   Step 1 — Frozen LTI baseline  ✓"),
    (0, "3.   Step 2 — LPV extension & ZOH validation  ✓"),
    (0, "4.   Step 3 — CT + RK4 approach decision"),
    (0, "5.   LPV-LFR structure"),
    (0, "6.   Open items & what we need from ASMPT"),
    (0, "7.   Timeline"),
]
add_bullet_block(slide, Inches(1.0), Inches(1.6), Inches(11.0), Inches(5.5),
                 items, font_size=20)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Project Context
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = WHITE
add_title_bar(slide, "Project Context", "Goal: augment the gantry FP model with data-driven corrections")
add_footer(slide, page=3)

# Two-column layout
y_start = Inches(1.55)
col_w   = Inches(5.9)
gap     = Inches(0.35)

# Left column header
section_box(slide, Inches(0.4), y_start, col_w, 0, "Physics Baseline (known)")
add_bullet_block(slide, Inches(0.5), y_start + Inches(0.48), col_w - Inches(0.1), Inches(3.2),
    [(0, "ASMPT dual-gantry FP model (García-Herreros et al.)"),
     (0, "Euler-Lagrange formulation"),
     (0, "3 DOF: X1, X2, Y axes"),
     (0, "CT ODE: M(Y)·ẍ + C·ẋ + K·x = P·u"),
     (0, "Y-dependent inertia → quasi-LPV"),
     (0, "Linearisation gap: Coriolis + Coulomb friction"),
    ], font_size=15)

# Right column header
rx = Inches(0.4) + col_w + gap
section_box(slide, rx, y_start, col_w, 0, "Augmentation Framework (learned)")
add_bullet_block(slide, rx + Inches(0.1), y_start + Inches(0.48), col_w - Inches(0.1), Inches(3.2),
    [(0, "LFR-based model augmentation (Hoekstra EJC 2025)"),
     (0, "Baseline + data-driven correction"),
     (0, "Interconnection matrix S couples physics + network"),
     (0, "Training on experimental gantry data"),
     (0, "Target: learn Coriolis + Coulomb gap"),
    ], font_size=15)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Step 1: Frozen LTI Baseline
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = WHITE
add_title_bar(slide, "Step 1 — Frozen LTI Baseline", "Python translation of MATLAB FP model")
add_footer(slide, page=4)
status_badge(slide, Inches(11.5), Inches(0.15), "✓ COMPLETE", done=True)

y = Inches(1.55)
section_box(slide, Inches(0.4), y, Inches(12.4), 0, "Implementation")
add_bullet_block(slide, Inches(0.5), y + Inches(0.48), Inches(12.0), Inches(1.8),
    [(0, "Y fixed at operating point Y = 0.3 m (linearisation point from main.m)"),
     (0, "ZOH discretisation at fs = 16 kHz  →  A, B, C, D matrices"),
     (0, "Stage coordinate transform via P matrix applied"),
    ], font_size=15)

y2 = y + Inches(2.3)
section_box(slide, Inches(0.4), y2, Inches(12.4), 0, "Validation Results")
add_bullet_block(slide, Inches(0.5), y2 + Inches(0.48), Inches(12.0), Inches(2.8),
    [(0, "Python A, B, C, D match MATLAB to < 1×10⁻¹⁰  (actual: ~1×10⁻¹⁹)  ✓"),
     (0, "Open-loop simulation vs Simscape (nonlinear ground truth):"),
     (1, "X1 residual: 4.86 µm RMS  —  X2: 3.19 µm  —  Y: 0.10 µm"),
     (1, "Bounded linearisation gap, dominant dynamics captured"),
     (0, "Linearisation gap (Coriolis + Coulomb) = augmentation target"),
    ], font_size=15)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Step 2: LPV Extension
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = WHITE
add_title_bar(slide, "Step 2 — LPV Extension", "Y-dependent inertia M(Y) as scheduling variable")
add_footer(slide, page=5)
status_badge(slide, Inches(11.5), Inches(0.15), "✓ COMPLETE", done=True)

y = Inches(1.55)
section_box(slide, Inches(0.4), y, Inches(12.4), 0, "Method: Frozen-at-sampling-instant ZOH  (Tóth 2010)")
add_bullet_block(slide, Inches(0.5), y + Inches(0.48), Inches(12.0), Inches(1.8),
    [(0, "At each sample k: evaluate A(Y_k), B(Y_k) via ZOH from CT model"),
     (0, "Justified at 16 kHz: ΔY ≤ 0.125 mm/sample  →  220:1 timescale separation"),
     (0, "Matrix exponential used: torch.linalg.matrix_exp  (differentiable, exact)"),
    ], font_size=15)

y2 = y + Inches(2.35)
section_box(slide, Inches(0.4), y2, Inches(12.4), 0, "Validation Results")
add_bullet_block(slide, Inches(0.5), y2 + Inches(0.48), Inches(12.0), Inches(2.8),
    [(0, "DT-LPV (Python) vs CT reference q1 (MATLAB ode45):"),
     (1, "BFR:  X1 = 99.99%    X2 = 99.98%    Y = 100.00%"),
     (1, "Sub-nanometre residual — ZOH discretisation confirmed correct"),
     (0, "What this proves: Python physics matches MATLAB to numerical precision"),
     (0, "Next: compare DT-LPV vs Simscape  →  defines augmentation target (Coriolis + Coulomb)"),
    ], font_size=15)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Step 3: CT + RK4 Approach Decision
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = WHITE
add_title_bar(slide, "Step 3 — Discretisation Approach Decision", "ZOH (Steps 1–2) → CT + RK4 (training loop)")
add_footer(slide, page=6)

# Left: original plan
lx = Inches(0.4)
y  = Inches(1.55)
section_box(slide, lx, y, Inches(5.9), 0, "Original Plan: Tóth (2010) ZOH for LPV-DT")
add_bullet_block(slide, lx + Inches(0.1), y + Inches(0.48), Inches(5.7), Inches(2.4),
    [(0, "Tóth (2010): 'Parameter-Varying Systems'"),
     (0, "Frozen-at-sampling-instant ZOH"),
     (0, "Gives exact DT LPV matrices A(Y_k), B(Y_k)"),
     (0, "Used and validated in Steps 1–2  ✓"),
     (0, "Pre-discretise before training loop"),
    ], font_size=14)

# Arrow
add_textbox(slide, Inches(6.4), Inches(2.5), Inches(0.8), Inches(0.5),
            "→", font_size=36, bold=True, color=ORANGE)

# Right: new decision
rx = Inches(7.1)
section_box(slide, rx, y, Inches(5.9), 0, "Supervisor Decision: Keep CT, Use RK4  (2026-03-20)")
add_bullet_block(slide, rx + Inches(0.1), y + Inches(0.48), Inches(5.7), Inches(2.4),
    [(0, "Keep model in continuous time"),
     (0, "RK4 with fixed step inside training loop"),
     (0, "ZOH held for validation (Steps 1–2) only"),
     (0, "No pre-discretisation needed"),
    ], font_size=14)

# Reasoning box
y3 = Inches(4.3)
section_box(slide, Inches(0.4), y3, Inches(12.4), 0, "Why CT + RK4 is preferred for the training loop", color=BLUE_DARK)
add_bullet_block(slide, Inches(0.5), y3 + Inches(0.48), Inches(12.0), Inches(2.5),
    [(0, "Steps 1–2 validated the physics with ZOH — that goal is achieved. No need to keep ZOH."),
     (0, "Pre-discretising with frozen ZOH before the LFR is constructed gets messy: "
          "scheduling enters via Δ(p), not via A(p) matrices in the training graph."),
     (0, "RK4 integrates the CT ODE directly at each time step — clean, no approximation error."),
     (0, "CT model keeps all physics intact; ZOH (zero-order hold on input) still assumed within each RK4 step."),
    ], font_size=14)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — LPV-LFR Structure
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = WHITE
add_title_bar(slide, "LPV-LFR Structure", "Linear Fractional Representation for LPV scheduling")
add_footer(slide, page=7)

y = Inches(1.55)
section_box(slide, Inches(0.4), y, Inches(12.4), 0, "LFR Definition  (Drenth 2025, Hoekstra 2025)")
add_bullet_block(slide, Inches(0.5), y + Inches(0.45), Inches(12.0), Inches(1.5),
    [(0, "LFR = pair  { G_b, Δ(p) }   where G_b is constant,  Δ(p) = diag(p·I_η)  is scheduling"),
     (0, "Scheduling variable p = Y  (gantry Y-position) — directly measured state"),
     (0, "Well-posedness: det(I − D_zw · Δ(p)) ≠ 0  for all p  →  enforced by D_zw = exp(−N), N > 0"),
    ], font_size=14)

y2 = y + Inches(2.1)
section_box(slide, Inches(0.4), y2, Inches(5.9), 0, "Baseline LFR")
add_bullet_block(slide, Inches(0.5), y2 + Inches(0.45), Inches(5.7), Inches(1.9),
    [(0, "Physics ODE:  M(Y)·ẍ + C·ẋ + K·x = P·u"),
     (0, "A_c(Y) has rational Y-entries  (via M(Y)⁻¹)"),
     (0, "LFT realization pulls Y-dependence into Δ(Y)"),
     (0, "Open: Zhou, Doyle & Glover (1996) Ch. 10"),
    ], font_size=14)

section_box(slide, Inches(6.85), y2, Inches(5.9), 0, "Augmentation LFR (learned)")
add_bullet_block(slide, Inches(6.95), y2 + Inches(0.45), Inches(5.7), Inches(1.9),
    [(0, "Neural network correction Δf(x, u)"),
     (0, "Wired via interconnection matrix S"),
     (0, "Learns Coriolis + Coulomb gap from data"),
     (0, "Ref: Hoekstra EJC Fig. 1 — block diagram"),
    ], font_size=14)

y3 = y2 + Inches(2.5)
section_box(slide, Inches(0.4), y3, Inches(12.4), 0, "M(Y) Invertibility  (required for LFR realization)")
add_bullet_block(slide, Inches(0.5), y3 + Inches(0.45), Inches(12.0), Inches(1.0),
    [(0, "Proved analytically: M(Y) is positive definite for ALL Y ∈ ℝ  (Sylvester's criterion)"),
     (1, "Leading minors: D₁ = 53.8  |  D₂ = 441Y²−3.66Y+210  |  D₃ = 4458Y²−37Y+2063"),
     (1, "All discriminants negative  →  no real roots  →  M(Y) invertible everywhere  ✓"),
    ], font_size=14)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Open Items & What We Need from ASMPT
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = WHITE
add_title_bar(slide, "Open Items & What We Need from ASMPT")
add_footer(slide, page=8)

lx = Inches(0.4)
y  = Inches(1.55)

section_box(slide, lx, y, Inches(5.9), 0, "Open Technical Items", color=ORANGE)
add_bullet_block(slide, lx + Inches(0.1), y + Inches(0.48), Inches(5.7), Inches(3.2),
    [(0, "Baseline LFR realization"),
     (1, "Converting rational M(Y)⁻¹ physics to LFR form"),
     (1, "Ref: Zhou, Doyle & Glover (1996) — being resolved"),
     (0, "Discretising LFRs"),
     (1, "Interaction between RK4 step and Δ(p) block"),
     (1, "Supervisor action item (2026-03-20): paper needed"),
     (0, "η choice: scheduling repetition count in Δ(Y)"),
     (0, "Sample rate: 16 kHz (measured) vs 20 kHz (spec)"),
    ], font_size=14)

rx = Inches(6.85)
section_box(slide, rx, y, Inches(5.9), 0, "Needed from ASMPT", color=BLUE_MID)
add_bullet_block(slide, rx + Inches(0.1), y + Inches(0.48), Inches(5.7), Inches(3.2),
    [(0, "Experimental measurement data"),
     (1, "Signals: [X1, X2, Y] position + applied forces [F_X1, F_X2, F_Y]"),
     (1, "Y must sweep operational range  (±350 mm)"),
     (0, "Excitation variety"),
     (1, "Step responses, PRBS, sinusoidal — ideally multiple"),
     (0, "Closed-loop controller (optional)"),
     (1, "Cfb matrix for full closed-loop simulation"),
     (0, "Confirmation: sample rate 16 kHz or 20 kHz?"),
    ], font_size=14)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Timeline
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = WHITE
add_title_bar(slide, "Timeline")
add_footer(slide, page=9)

rows = [
    ("Step 1", "Frozen LTI Baseline",         "Feb 2026",     True),
    ("Step 2", "LPV Extension & ZOH Validation","Mar 2026",    True),
    ("Step 3", "CT + RK4 + LFR Construction",  "Mar–Apr 2026", False),
    ("Step 4", "Three Research Novelties",      "Apr–May 2026", False),
    ("—",      "April 9 Supervisor Meeting",    "09 Apr 2026",  False),
    ("—",      "ASMPT Data Integration & Training", "May 2026", False),
]

y = Inches(1.65)
for step, desc, date, done in rows:
    color = GREEN_DARK if done else BLUE_MID
    # Step pill
    add_rect(slide, Inches(0.4), y, Inches(1.1), Inches(0.42), fill_color=color)
    add_textbox(slide, Inches(0.42), y + Inches(0.04), Inches(1.06), Inches(0.34),
                step, font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Description
    add_textbox(slide, Inches(1.65), y, Inches(8.5), Inches(0.42),
                desc, font_size=14, color=DARK_GREY)
    # Date
    add_textbox(slide, Inches(10.3), y, Inches(2.6), Inches(0.42),
                date, font_size=13, color=RGBColor(0x55, 0x55, 0x55))
    # Status tick
    tick = "✓" if done else "→"
    add_textbox(slide, Inches(12.95), y, Inches(0.35), Inches(0.42),
                tick, font_size=16, bold=True,
                color=GREEN_DARK if done else ORANGE)
    y += Inches(0.62)


# ════════════════════════════════════════════════════════════════════════════
# Save
# ════════════════════════════════════════════════════════════════════════════
os.makedirs("presentations", exist_ok=True)
out_path = "presentations/asmpt-23-03-2026.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
